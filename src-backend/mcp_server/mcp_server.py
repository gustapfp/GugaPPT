import json
import os
from datetime import datetime
from typing import Literal

import numpy as np
from matplotlib import pyplot as plt
from mcp.server.fastmcp import FastMCP
from pptx import Presentation
from pptx.util import Inches
from tavily import TavilyClient  # type: ignore

from core.consts import (
    BODY_FONT_SIZE,
    BODY_FONT_SIZE_WITH_IMAGE,
    BODY_LINE_SPACING,
    BODY_WIDTH_WITH_IMAGE,
    DOMAIN_BLACKLIST,
    FILE_PATH,
    IMAGE_HEIGHT,
    SLIDE_HEIGHT,
    SLIDE_WIDTH,
)
from core.logger_config import logger
from core.settings import settings
from mcp_server.helper.ppt_style import apply_body_style, apply_title_style
from mcp_server.helper.source_validator import source_validator

mcp_server = FastMCP("PPT-Generator-Tools")

tavily_client = TavilyClient(api_key=settings.TAVILY_API_KEY)


@mcp_server.tool(
    name="search_web",
    description="Search the web for information",
)
def search_web(query: str, search_depth: Literal["basic", "advanced"] = "advanced") -> str:
    """Search the web for information based on the given query.

    Args:
        query (str): The query to search the web for.
        search_depth (Literal["basic", "advanced"]): The depth of the search.
    Returns:
        str: The information searched for.
    """
    logger.info("Searching the web for information...")
    try:
        response = tavily_client.search(
            query=query,
            search_depth=search_depth,
            max_results=10,
            exclude_domains=DOMAIN_BLACKLIST,
            chunks_per_source=3,
        )

        context = [{"content": r["content"], "url": r["url"]} for r in response.get("results", [])]
        logger.info(f"Context: {context}")
        ranked_results = source_validator.rank_sources(context)
        high_quality_results = [
            res for res in ranked_results if res["validation"]["tier"] in ["S", "A"]
        ]
        if not high_quality_results:
            logger.warning(f"No Tier S/A results found for '{query}'. Returning empty list.")
            return json.dumps([])
        logger.info(f"Returning {len(high_quality_results)} high-quality results for '{query}'")
        return json.dumps(high_quality_results, indent=2)
    except Exception as e:
        return f"Error searching web: {str(e)}"


@mcp_server.tool(
    name="create_presentation",
    description="Create a PowerPoint presentation based on the given slides content.",
)
def create_presentation(filename: str, slides_content: str) -> str:
    """Create a PowerPoint presentation based on the given slides content.

    Args:
        filename (str): The filename of the presentation.
        slides_content (str): The slides content of the presentation.

    Returns:
        str: The message indicating the success or failure of the operation.
    """
    try:
        data = json.loads(slides_content)

        prs = Presentation()

        for slide_data in data:
            slide_layout = prs.slide_layouts[0]
            for layout in prs.slide_layouts:
                layout_name = layout.name.lower()
                if "content" in layout_name or "text" in layout_name:
                    slide_layout = layout
                    break
                if "title" in layout_name and "only" not in layout_name:
                    slide_layout = layout

            slide = prs.slides.add_slide(slide_layout)

            image_path = slide_data.get("image")
            has_image = image_path and os.path.exists(image_path)

            # -- Title --
            title = slide.shapes.title
            if title:
                title.text = slide_data.get("title", "No Title")
                apply_title_style(title)

            # -- Body --

            body_shape = None
            for shape in slide.placeholders:
                if hasattr(shape, "placeholder_format") and shape.placeholder_format.idx == 1:
                    body_shape = shape
                    break

            if body_shape is None:
                placeholders = list(slide.placeholders)
                if len(placeholders) > 1:
                    body_shape = placeholders[1]

            if body_shape:
                tf = body_shape.text_frame  # pyright: ignore[reportAttributeAccessIssue]
                tf.word_wrap = True

                if has_image:
                    body_shape.width = BODY_WIDTH_WITH_IMAGE

                font_size = BODY_FONT_SIZE_WITH_IMAGE if has_image else BODY_FONT_SIZE

                points = slide_data.get("points", [])

                if points:
                    tf.paragraphs[0].text = points[0]
                    apply_body_style(tf.paragraphs[0], font_size)
                    tf.paragraphs[0].level = 0

                    for point in points[1:]:
                        p = tf.add_paragraph()
                        p.text = point
                        apply_body_style(p, font_size)
                        p.level = 0
                        p.space_before = BODY_LINE_SPACING

            # -- Speaker Notes & Sources --
            speaker_notes = slide_data.get("speaker_notes", "")
            sources = slide_data.get("sources", [])
            if speaker_notes or sources:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                content = speaker_notes if speaker_notes else ""

                if sources:
                    if content:
                        content += "\n\n"
                    content += "Sources:\n" + "\n".join(f"- {url}" for url in sources)
                if text_frame:
                    text_frame.text = content

            # -- Image --
            if has_image:
                try:
                    picture = slide.shapes.add_picture(
                        image_path,
                        left=Inches(0),
                        top=Inches(0),
                        height=IMAGE_HEIGHT,
                    )
                    picture.left = (SLIDE_WIDTH - picture.width) // 2
                    picture.top = (SLIDE_HEIGHT - picture.height) // 2
                except Exception as e:
                    logger.warning(f"Could not add image {image_path}: {e}")

        path = FILE_PATH / f"{filename}.pptx"
        path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(path.resolve()))
        return f"Successfully saved presentation to {path}"
    except Exception as e:
        return f"Error creating PPT: {str(e)}"


@mcp_server.tool(
    name="generate_chart",
    description="Generate visual assets for the presentation.",
)
def generate_chart(data_json: str, chart_type: str, title: str) -> str:
    """
    Generates a statistical chart (bar, pie, or line) and saves it as a PNG image.

    Args:
        data_json: A JSON string containing 'labels' (list) and 'values' (list).
                   Example: '{"labels": ["Q1", "Q2"], "values": [100, 150]}'
        chart_type: The type of chart to generate. Options: "bar", "pie", "line".
        title: The title of the chart.

    Returns:
        The file path of the generated image.
    """
    logger.info(
        f"Generate Visual Assets Tool was triggered with data_json: {data_json}, chart_type: {chart_type}, title: {title}"
    )
    try:
        data = json.loads(data_json)
        labels = data.get("labels", [])
        values = data.get("values", [])
        unit = data.get("unit", "Values")

        if not labels or not values:
            raise ValueError("Error: JSON must contain 'labels' and 'values' lists.")

        if len(labels) != len(values):
            raise ValueError("Error: 'labels' and 'values' must have the same length.")

        plt.figure(figsize=(10, 6))

        if chart_type.lower() == "bar":
            plt.bar(labels, values, color="#4F81BD")
            plt.xlabel("Categories")
            plt.ylabel(unit)

        elif chart_type.lower() == "line":
            plt.plot(labels, values, marker="o", linestyle="-", color="#C0504D", linewidth=2)
            plt.ylabel(unit)
            plt.grid(True, linestyle="--", alpha=0.7)

        elif chart_type.lower() == "pie":
            cmap = plt.get_cmap("Paired")
            rgba = cmap(np.linspace(0, 1, len(values)))
            colors = [tuple(rgba[i]) for i in range(len(values))]
            plt.pie(
                values,
                labels=labels,
                autopct="%1.0f%%",
                startangle=90,
                colors=colors,
            )
            plt.ylabel(unit)
        else:
            return f"Error: Unsupported chart type '{chart_type}'. Use 'bar', 'pie', or 'line'."

        plt.title(title)

        safe_title = "".join(c if c.isalnum() else "_" for c in title)
        timestamp = int(datetime.now().timestamp())
        filename = f"chart_{safe_title}_{timestamp}.png"
        path = FILE_PATH / "charts" / filename
        path.parent.mkdir(parents=True, exist_ok=True)
        plt.savefig(path, bbox_inches="tight", dpi=100)
        plt.close()

        print(f"DEBUG: Saved chart to {path}")
        return str(path)

    except json.JSONDecodeError:
        return "Error: Invalid JSON string provided."
    except Exception as e:
        return f"Error generating chart: {str(e)}"


if __name__ == "__main__":
    mcp_server.run()
