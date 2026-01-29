from typing import Literal
from mcp_server.mcp_server import mcp_server
from structlog import get_logger
from pptx import Presentation
import json
from core.consts import DOMAIN_BLACKLIST, FILE_PATH
from tavily import TavilyClient  # type: ignore
from core.settings import settings

logger = get_logger()


tavily_client = TavilyClient(api_key=settings.TAVILY_API_KEY)


@mcp_server.tool(
    name="search_web",
    description="Search the web for information",
)
def search_web(query: str, search_depth: Literal["basic", "advanced"] = "basic") -> str:
    """Search the web for information based on the given query.

    Args:
        query (str): The query to search the web for.
        search_depth (Literal["basic", "advanced"]): The depth of the search.
    Returns:
        str: The information searched for.
    """
    logger.info(f"Search Web Tool was triggered with query: {query}. Searching the web for information...")
    try:
        response = tavily_client.search(
            query=query,
            search_depth=search_depth,
            max_results=3,
            exclude_domains=DOMAIN_BLACKLIST,
        )

        context = [{"content": r["content"], "url": r["url"]} for r in response.get("results", [])]
        return json.dumps(context)
    except Exception as e:
        return f"Error searching web: {str(e)}"


@mcp_server.tool(
    name="create_presentation",
    description="Create a PowerPoint presentation based on the given slides content.",
)
def create_presentation(filename: str, slides_content: str) -> str:
    """Create a PowerPoint presentation based on the given slides content.

    Args:
        filename (str): The filename of the presentation.
        slides_content (str): The slides content of the presentation.

    Returns:
        str: The message indicating the success or failure of the operation.
    """
    try:
        data = json.loads(slides_content)
        prs = Presentation()

        for slide_data in data:
            # -- Bullet layout --
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            # -- Title --
            title = slide.shapes.title
            if title:
                title.text = slide_data.get("title", "No Title")

            # -- Body --
            body_shape = slide.placeholders[1]
            tf = getattr(body_shape, "text_frame")
            for point in slide_data.get("points", []):
                p = tf.add_paragraph()
                p.text = point

        path = FILE_PATH / f"{filename}.pptx"
        path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(path.resolve()))
        return f"Successfully saved presentation to {path}"
    except Exception as e:
        return f"Error creating PPT: {str(e)}"
