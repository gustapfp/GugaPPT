import json
import os
import tempfile
from pathlib import Path
from unittest.mock import AsyncMock, MagicMock, patch

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation as PptxPresentation


class TestSourceValidator:
    """Tests for SourceValidator helper."""

    def test_normalize_url_removes_query_params(self):
        """Test URL normalization removes query parameters."""
        from mcp_server.helper.source_validator import SourceValidator

        validator = SourceValidator()
        url = "https://example.com/article?utm_source=google&ref=123"
        normalized = validator.normalize_url(url)
        assert normalized == "https://example.com/article"

    def test_get_metadata_extracts_author(self):
        """Test metadata extraction for author."""
        from bs4 import BeautifulSoup

        from mcp_server.helper.source_validator import SourceValidator

        validator = SourceValidator()
        html = '<html><head><meta name="author" content="John Doe"></head></html>'
        soup = BeautifulSoup(html, "html.parser")
        meta = validator.get_metadata(soup)
        assert meta["author"] == "John Doe"

    def test_get_metadata_extracts_date(self):
        """Test metadata extraction for date."""
        from bs4 import BeautifulSoup

        from mcp_server.helper.source_validator import SourceValidator

        validator = SourceValidator()
        html = '<html><head><meta name="date" content="2026-01-30"></head><body><h2>References</h2><p>Source list</p></body></html>'
        soup = BeautifulSoup(html, "html.parser")
        meta = validator.get_metadata(soup)
        assert meta["date"] == "2026-01-30"
        assert meta["has_references"] is True

    @patch("mcp_server.helper.source_validator.requests.get")
    def test_validate_url_live_site(self, mock_get):
        """Test URL validation for a live site."""
        from mcp_server.helper.source_validator import SourceValidator

        mock_response = MagicMock()
        mock_response.status_code = 200
        mock_response.content = b"<html><head><meta name='author' content='Test'></head></html>"
        mock_get.return_value = mock_response

        validator = SourceValidator()
        result = validator.validate_url("https://example.com", tavily_confidence=0.8)

        assert result["status"] == "live"
        assert result["score"] > 0
        assert result["tier"] in ["S", "A", "B"]

    @patch("mcp_server.helper.source_validator.requests.get")
    def test_validate_url_dead_site(self, mock_get):
        """Test URL validation for a dead site."""
        from mcp_server.helper.source_validator import SourceValidator

        mock_response = MagicMock()
        mock_response.status_code = 404
        mock_get.return_value = mock_response

        validator = SourceValidator()
        result = validator.validate_url("https://example.com/404", tavily_confidence=0.8)

        assert result["status"] == "dead"
        assert result["tier"] == "C"


class TestPlannerAgent:
    """Tests for PlannerAgent."""

    @pytest.mark.asyncio
    async def test_create_presentation_plan_success(self):
        """Test successful presentation plan creation."""
        from mcp_server.agents.planner.agent import PlannerAgent
        from mcp_server.agents.planner.schemas import (
            PresentationPayload,
            PresentationPlan,
            SlidePlan,
        )

        agent = PlannerAgent()

        mock_plan = PresentationPlan(
            topic="Test Topic",
            slides=[
                SlidePlan(
                    slide_number=i,
                    title=f"Slide {i}",
                    search_queries=[f"query{i}"],
                    content_goal=f"Goal {i}",
                )
                for i in range(3)
            ],
        )

        mock_response = MagicMock()
        mock_response.choices = [MagicMock(message=MagicMock(parsed=mock_plan))]

        with patch.object(
            agent.client.beta.chat.completions, "parse", new_callable=AsyncMock
        ) as mock_parse:
            mock_parse.return_value = mock_response

            payload = PresentationPayload(topic="Test Topic", num_slides=3)
            result = await agent.create_presentation_plan(payload)

            assert result.topic == "Test Topic"
            assert len(result.slides) == 3
            mock_parse.assert_called_once()

    @pytest.mark.asyncio
    async def test_validate_response_retries_on_none(self):
        """Test validation retries when response is None."""
        from mcp_server.agents.planner.agent import PlannerAgent
        from mcp_server.agents.planner.schemas import (
            PresentationPayload,
            PresentationPlan,
            SlidePlan,
        )

        agent = PlannerAgent()
        payload = PresentationPayload(topic="Test", num_slides=2)

        valid_plan = PresentationPlan(
            topic="Test",
            slides=[
                SlidePlan(
                    slide_number=i,
                    title=f"Slide {i}",
                    search_queries=["q"],
                    content_goal="Goal",
                )
                for i in range(2)
            ],
        )

        mock_responses = [
            MagicMock(choices=[MagicMock(message=MagicMock(parsed=None))]),
            MagicMock(choices=[MagicMock(message=MagicMock(parsed=valid_plan))]),
        ]

        with patch.object(
            agent.client.beta.chat.completions, "parse", new_callable=AsyncMock
        ) as mock_parse:
            mock_parse.side_effect = mock_responses

            result = await agent.create_presentation_plan(payload)
            assert result.topic == "Test"
            assert mock_parse.call_count == 2


class TestResearcherAgent:
    """Tests for ResearcherAgent."""

    @pytest.mark.asyncio
    async def test_research_web_empty_results(self):
        """Test research_web handles empty results."""
        from mcp_server.agents.researcher.agent import ResearcherAgent
        from mcp_server.agents.researcher.schemas import ResearcherPayload

        agent = ResearcherAgent()
        mock_session = AsyncMock()

        mock_result = MagicMock()
        mock_result.content = []
        mock_session.call_tool.return_value = mock_result

        payload = ResearcherPayload(slide_title="Test Slide", search_queries=["test query"])

        result = await agent.research_web(payload, mock_session)

        assert result.slide_topic == "Test Slide"
        assert len(result.facts) == 0

    @pytest.mark.asyncio
    async def test_summarize_facts_success(self):
        """Test successful fact summarization."""
        from mcp_server.agents.researcher.agent import ResearcherAgent
        from mcp_server.agents.researcher.schemas import Fact, ResearchSummary

        agent = ResearcherAgent()

        mock_summary = ResearchSummary(
            slide_topic="AI Trends",
            facts=[
                Fact(content="AI is growing fast", source_url="https://example.com"),
            ],
        )

        mock_response = MagicMock()
        mock_response.choices = [MagicMock(message=MagicMock(parsed=mock_summary))]

        with patch.object(
            agent.client.beta.chat.completions, "parse", new_callable=AsyncMock
        ) as mock_parse:
            mock_parse.return_value = mock_response

            result = await agent.summarize_facts(
                raw_context=["Some context"], slide_title="AI Trends"
            )

            assert result.slide_topic == "AI Trends"
            assert len(result.facts) == 1


class TestWriterAgent:
    """Tests for WriterAgent."""

    @pytest.mark.asyncio
    async def test_prepare_presentation_success(self):
        """Test successful presentation preparation."""
        from mcp_server.agents.writer.agent import WriterAgent
        from mcp_server.agents.writer.schemas import (
            ChartData,
            PresentationContent,
            SlideContent,
            VisualRequest,
        )

        agent = WriterAgent()

        mock_content = PresentationContent(
            filename_suggestion="test_presentation",
            slides=[
                SlideContent(
                    title="Introduction",
                    points=["Point 1", "Point 2"],
                    speaker_notes="Notes",
                    sources=["https://source.com"],
                    visual_request=VisualRequest(
                        type="chart",
                        prompt="Revenue Chart",
                        data_json=ChartData(labels=["Q1", "Q2"], values=[100.0, 200.0], unit="USD"),
                    ),
                ),
            ],
        )

        mock_response = MagicMock()
        mock_response.choices = [MagicMock(message=MagicMock(parsed=mock_content))]

        with patch.object(
            agent.client.beta.chat.completions, "parse", new_callable=AsyncMock
        ) as mock_parse:
            mock_parse.return_value = mock_response

            result = await agent.prepare_presentation(
                topic="Test",
                plan_json={"topic": "Test", "slides": []},
                research_data=[],
            )

            assert result.filename_suggestion == "test_presentation"
            assert len(result.slides) == 1

    @pytest.mark.asyncio
    async def test_write_presentation_calls_mcp_tool(self):
        """Test write_presentation calls the MCP tool correctly."""
        from mcp_server.agents.writer.agent import WriterAgent
        from mcp_server.agents.writer.schemas import (
            PresentationContent,
            SlideContent,
        )

        agent = WriterAgent()
        mock_session = AsyncMock()

        content = PresentationContent(
            filename_suggestion="test",
            slides=[
                SlideContent(
                    title="Test Slide",
                    points=["Point 1"],
                    speaker_notes=None,
                    sources=None,
                ),
            ],
        )

        await agent.write_presentation(content=content, session=mock_session, filename="test_file")

        mock_session.call_tool.assert_called_once_with(
            "create_presentation",
            arguments={
                "filename": "test_file",
                "slides_content": json.dumps(
                    [
                        {
                            "title": "Test Slide",
                            "points": ["Point 1"],
                            "speaker_notes": None,
                            "sources": None,
                        }
                    ]
                ),
            },
        )

    @pytest.mark.asyncio
    async def test_validate_response_requires_chart(self):
        """Test validation requires at least one chart."""
        from mcp_server.agents.writer.agent import WriterAgent
        from mcp_server.agents.writer.schemas import (
            PresentationContent,
            SlideContent,
        )

        agent = WriterAgent()
        agent.retry_count = 3

        content = PresentationContent(
            filename_suggestion="test",
            slides=[
                SlideContent(
                    title="Test",
                    points=["Point"],
                    speaker_notes=None,
                    sources=None,
                    visual_request=None,
                ),
            ],
        )

        with pytest.raises(ValueError, match="No chart generated"):
            await agent._validate_response(content, "Test", {}, [])


class TestIllustratorAgent:
    """Tests for IllustratorAgent."""

    @pytest.mark.asyncio
    async def test_create_visuals_chart(self):
        """Test creating chart visuals."""
        from mcp_server.agents.illustrator.agent import IllustratorAgent

        agent = IllustratorAgent()
        mock_session = AsyncMock()

        mock_result = MagicMock()
        mock_result.content = [MagicMock(text="/path/to/chart.png")]
        mock_session.call_tool.return_value = mock_result

        visual_requests = [
            {
                "slide_number": 0,
                "type": "chart",
                "prompt": "Revenue Chart",
                "data_json": {"labels": ["Q1", "Q2"], "values": [100, 200], "unit": "USD"},
            }
        ]

        result = await agent.create_visuals(visual_requests, mock_session)

        assert len(result.assets) == 1
        assert result.assets[0].slide_number == 0
        assert result.assets[0].asset_type == "chart"
        assert result.assets[0].file_path == "/path/to/chart.png"

    @pytest.mark.asyncio
    async def test_create_visuals_handles_exception(self):
        """Test create_visuals handles exceptions gracefully."""
        from mcp_server.agents.illustrator.agent import IllustratorAgent

        agent = IllustratorAgent()
        mock_session = AsyncMock()
        mock_session.call_tool.side_effect = Exception("Tool error")

        visual_requests = [
            {
                "slide_number": 0,
                "type": "chart",
                "prompt": "Chart",
                "data_json": {"labels": ["A"], "values": [1], "unit": "X"},
            }
        ]

        result = await agent.create_visuals(visual_requests, mock_session)

        assert len(result.assets) == 0


class TestMcpServerTools:
    """Tests for MCP server tools."""

    @pytest.mark.parametrize("chart_type", ["bar", "line", "pie"])
    def test_generate_chart_valid_types(self, chart_type):
        """Test chart generation for all valid types."""
        from mcp_server.mcp_server import generate_chart

        with (
            tempfile.TemporaryDirectory() as tmpdir,
            patch("mcp_server.mcp_server.FILE_PATH", Path(tmpdir)),
        ):
            data_json = json.dumps({"labels": ["A", "B", "C"], "values": [10, 20, 30], "unit": "X"})
            result = generate_chart(data_json, chart_type, f"Test {chart_type}")
            assert "chart_" in result and result.endswith(".png")

    def test_create_presentation(self):
        """Test presentation creation (success and error)."""
        from mcp_server.mcp_server import create_presentation

        assert "Error" in create_presentation("test", "not valid json")

        with (
            tempfile.TemporaryDirectory() as tmpdir,
            patch("mcp_server.mcp_server.FILE_PATH", Path(tmpdir)),
        ):
            slides = json.dumps(
                [{"title": "Slide", "points": ["Point"], "speaker_notes": None, "sources": None}]
            )
            result = create_presentation("test_ppt", slides)
            assert "Successfully saved" in result
            assert os.path.exists(Path(tmpdir) / "test_ppt.pptx")

    @patch("mcp_server.mcp_server.tavily_client")
    @patch("mcp_server.mcp_server.source_validator")
    def test_search_web(self, mock_validator, mock_tavily):
        """Test web search with various scenarios."""
        from mcp_server.mcp_server import search_web

        mock_tavily.search.return_value = {
            "results": [{"content": "Content", "url": "https://example.com"}]
        }
        mock_validator.rank_sources.return_value = [
            {
                "content": "Content",
                "url": "https://example.com",
                "validation": {"tier": "S", "score": 90},
            }
        ]
        result = json.loads(search_web("query"))
        assert len(result) == 1 and result[0]["validation"]["tier"] == "S"

        mock_validator.rank_sources.return_value = [
            {"content": "Low", "url": "https://bad.com", "validation": {"tier": "C", "score": 30}}
        ]
        assert json.loads(search_web("query")) == []

        mock_tavily.search.side_effect = Exception("API Error")
        assert "Error" in search_web("query")


class TestPresentationRoutes:
    """Tests for presentation API routes."""

    @pytest.fixture
    def client(self):
        """Create test client."""
        from fastapi import FastAPI

        from app.routes.presentation.router import presentation_router

        app = FastAPI()
        app.include_router(presentation_router)
        return TestClient(app)

    def test_generate_ppt_success(self, client):
        """Test successful presentation generation."""
        with patch("app.routes.presentation.router.run_ppt_workflow"):
            response = client.post(
                "/presentation/generate_ppt",
                json={"topic": "AI Trends", "slides": 5},
            )

            assert response.status_code == 202
            data = response.json()
            assert data["status"] == "Success"
            assert "pprt_id" in data
            assert "AI_Trends" in data["pprt_id"]

    def test_download_ppt_found(self, client):
        """Test downloading existing presentation."""
        with tempfile.TemporaryDirectory() as tmpdir:
            test_path = Path(tmpdir) / "test-123.pptx"
            prs = PptxPresentation()
            prs.save(str(test_path))

            with patch("app.routes.presentation.router.FILE_PATH", Path(tmpdir)):
                response = client.get("/presentation/download/test-123")

                assert response.status_code == 200
                assert (
                    response.headers["content-type"]
                    == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )


if __name__ == "__main__":
    pytest.main([__file__, "-v", "--tb=short"])
